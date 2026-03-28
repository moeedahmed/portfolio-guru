"""
Document extraction for Portfolio Guru — handles PDF, PPTX, DOCX files.
Uses markitdown for PPTX/DOCX and pypdf/pdfplumber for PDFs.
"""
import os
import tempfile
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


async def extract_from_document(file_path: str) -> str:
    """
    Extract text from a document file (PDF, PPTX, DOCX).
    Returns extracted text or empty string if extraction fails.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()
    
    try:
        if suffix == ".pdf":
            return await _extract_pdf(file_path)
        elif suffix in (".pptx", ".ppt"):
            return await _extract_pptx(file_path)
        elif suffix in (".docx", ".doc"):
            return await _extract_docx(file_path)
        elif suffix in (".txt", ".md", ".rst"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        else:
            logger.warning(f"Unsupported file type: {suffix}")
            return ""
    except Exception as e:
        logger.error(f"Document extraction failed for {file_path}: {e}")
        return ""


async def _ocr_pdf_with_gemini(file_path: str) -> str:
    """OCR a scanned PDF using Gemini Vision (native PDF input — no page rendering needed)."""
    import asyncio
    from google import genai
    from google.genai import types

    client = genai.Client(api_key=os.environ.get("GOOGLE_API_KEY"))

    with open(file_path, "rb") as f:
        pdf_data = f.read()

    prompt = (
        "Extract ALL text from this PDF document exactly as written. "
        "Preserve the structure (headings, paragraphs, lists, tables). "
        "Return only the extracted text, no commentary."
    )

    contents = [
        types.Content(
            role="user",
            parts=[
                types.Part.from_bytes(data=pdf_data, mime_type="application/pdf"),
                types.Part.from_text(text=prompt),
            ]
        )
    ]

    loop = asyncio.get_event_loop()
    models_to_try = ["gemini-3-flash-preview", "gemini-2.5-flash"]

    for model in models_to_try:
        try:
            response = await loop.run_in_executor(
                None,
                lambda m=model: client.models.generate_content(model=m, contents=contents)
            )
            return response.text.strip()
        except Exception as e:
            error_msg = str(e).lower()
            if any(t in error_msg for t in ["503", "unavailable", "overloaded", "404", "429", "not found"]):
                logger.warning(f"Gemini {model} PDF OCR failed: {e} — trying next model")
                continue
            raise

    logger.error("All Gemini models failed for PDF OCR — returning empty string")
    return ""


async def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF. Falls back to Gemini Vision OCR for scanned/image-only PDFs."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        combined = "\n\n".join(text_parts)

        if len(combined.strip()) > 200:
            return combined

        logger.info("PDF returned minimal text from pypdf — falling back to Gemini Vision OCR")
    except Exception as e:
        logger.warning(f"pypdf extraction failed: {e} — trying Gemini Vision OCR")

    return await _ocr_pdf_with_gemini(file_path)


async def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return ""


async def _extract_docx(file_path: str) -> str:
    """Extract text from Word document using python-docx."""
    try:
        from docx import Document
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text.strip())
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return ""


def get_supported_extensions() -> list:
    """Return list of supported file extensions."""
    return [".pdf", ".pptx", ".ppt", ".docx", ".doc", ".txt", ".md"]


def is_supported_document(filename: str) -> bool:
    """Check if a filename has a supported extension."""
    suffix = Path(filename).suffix.lower()
    return suffix in get_supported_extensions()
